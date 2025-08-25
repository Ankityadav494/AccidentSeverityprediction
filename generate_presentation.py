#!/usr/bin/env python3
"""
Generate a PowerPoint presentation for the Road Accident Severity project.
Outputs: Road_Accident_Severity_Presentation.pptx
"""

from pathlib import Path
import sys
import pandas as pd

BASE_DIR = Path(__file__).resolve().parent
SRC_DIR = BASE_DIR / "src"
if str(SRC_DIR) not in sys.path:
    sys.path.insert(0, str(SRC_DIR))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from simple_prediction import SimpleAccidentPredictor


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        p.text = bullet
        p.level = 0


def add_table(prs: Presentation, title: str, df: pd.DataFrame) -> None:
    slide_layout = prs.slide_layouts[5]  # title only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    rows, cols = df.shape
    left = Inches(0.7)
    top = Inches(1.8)
    width = Inches(8.0)
    height = Inches(0.8 + 0.3 * (rows + 1))

    table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table
    # header
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        cell.text_frame.paragraphs[0].font.bold = True
    # data
    for i in range(rows):
        for j in range(cols):
            table.cell(i + 1, j).text = str(df.iat[i, j])


def add_feature_importance_slide(prs: Presentation, predictor: SimpleAccidentPredictor, model_name: str) -> None:
    fi = predictor.get_feature_importance(model_name=model_name)
    if fi is None or fi.empty:
        return
    topn = fi.sort_values("importance", ascending=False).head(10)
    add_table(prs, f"Top Feature Importance ({model_name})", topn)


def build_presentation(output_path: Path) -> None:
    prs = Presentation()

    add_title_slide(
        prs,
        title="Road Accident Severity Prediction",
        subtitle="Simple ML pipeline with Streamlit app",
    )

    add_bullets_slide(
        prs,
        title="Objective",
        bullets=[
            "Predict accident severity (Minor/Moderate/Major/Fatal)",
            "Use 10 essential features: time, location, road, weather",
            "Deliver interactive Streamlit app for scenario exploration",
        ],
    )

    add_bullets_slide(
        prs,
        title="Data & Features",
        bullets=[
            "Source: road_accident_dataset.csv",
            "Categoricals encoded (Country, Month, Day, Time, Urban/Rural, Road Type, Weather)",
            "Numericals scaled (Year, Visibility Level, Number of Vehicles)",
        ],
    )

    add_bullets_slide(
        prs,
        title="Models",
        bullets=[
            "DecisionTreeClassifier",
            "RandomForestClassifier (default)",
            "Saved artifacts: models/simple_*.pkl",
        ],
    )

    predictor = SimpleAccidentPredictor(models_dir=str(BASE_DIR / "models"))
    add_feature_importance_slide(prs, predictor, model_name="random_forest")

    add_bullets_slide(
        prs,
        title="Evaluation",
        bullets=[
            "Metrics computed during training (accuracy + CV)",
            "Confusion matrices saved under results/ if training was run",
            "Feature importance available for tree-based models",
        ],
    )

    add_bullets_slide(
        prs,
        title="App Demo (Streamlit)",
        bullets=[
            "Interactive inputs + model selection",
            "Prediction with confidence bars",
            "Sample scenarios and quick reset",
        ],
    )

    add_bullets_slide(
        prs,
        title="Next Steps",
        bullets=[
            "Add more robust handling for unseen categories",
            "Calibrate probabilities; try gradient boosting",
            "Collect more data and perform hyperparameter tuning",
        ],
    )

    prs.save(str(output_path))


if __name__ == "__main__":
    out = BASE_DIR / "Road_Accident_Severity_Presentation.pptx"
    build_presentation(out)
    print(f"Saved presentation to {out}")



